import PyPDF2
import io
import docx
from pptx import Presentation

def extract_text_from_pdf(file_stream) -> str:
    """Extracts all text from a PDF file stream with page markers."""
    try:
        reader = PyPDF2.PdfReader(file_stream)
        text = ""
        for idx, page in enumerate(reader.pages):
            extracted = page.extract_text()
            if extracted:
                text += f"\n--- [Page {idx + 1}] ---\n" + extracted + "\n"
        return text
    except Exception as e:
        print(f"Error extracting PDF: {e}")
        return ""

def extract_text_from_docx(file_stream) -> str:
    """Extracts all text from a DOCX file stream with paragraph markers."""
    try:
        doc = docx.Document(file_stream)
        text = ""
        for idx, para in enumerate(doc.paragraphs):
            if para.text.strip():
                text += f"[Para {idx + 1}] {para.text}\n"
        return text
    except Exception as e:
        print(f"Error extracting DOCX: {e}")
        return ""

def extract_text_from_pptx(file_stream) -> str:
    """Extracts all text from a PPTX file stream with slide markers."""
    try:
        prs = Presentation(file_stream)
        text = ""
        for idx, slide in enumerate(prs.slides):
            slide_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text += shape.text + "\n"
            if slide_text.strip():
                text += f"\n--- [Slide {idx + 1}] ---\n" + slide_text
        return text
    except Exception as e:
        print(f"Error extracting PPTX: {e}")
        return ""

def process_uploaded_files(uploaded_files):
    """
    Process a list of Streamlit uploaded files.
    Returns a tuple of (combined_text_string, list_of_filenames)
    """
    combined_text = ""
    file_names = []
    
    for file in uploaded_files:
        text = ""
        if file.name.lower().endswith('.pdf'):
            text = extract_text_from_pdf(file)
        elif file.name.lower().endswith('.docx'):
            text = extract_text_from_docx(file)
        elif file.name.lower().endswith('.pptx'):
            text = extract_text_from_pptx(file)
            
        if text.strip():
            combined_text += f"\n\n=== Start of Document: {file.name} ===\n{text}\n=== End of Document: {file.name} ===\n\n"
            file_names.append(file.name)
            
    return combined_text, file_names
